import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Define paths
TEMPLATE_PATH = "/home/elissa/Documents/ALU/CAPSTON/defence slides /Elissa_Twizeyimana Mission Capstone Defense Template.pptx"
OUTPUT_PATH = "/home/elissa/Documents/ALU/CAPSTON/defence slides /Inzira_Capstone_Defense.pptx"

# Diagrams
SYS_ARCH_PATH = "/home/elissa/Documents/ALU/CAPSTON/Inzira/docs/diagrams/fig2_system_architecture.png"
ERD_PATH = "/home/elissa/Documents/ALU/CAPSTON/Inzira/docs/diagrams/fig4_er_diagram.png"
ML_FLOW_PATH = "/home/elissa/Documents/ALU/CAPSTON/Inzira/docs/diagrams/fig1_ml_model_flow.png"

# Load presentation
prs = Presentation(TEMPLATE_PATH)

print(f"Loaded template from {TEMPLATE_PATH}")

# Helper to do paragraph text replacements
def replace_text_in_paragraphs(slide, replacements):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for old, new in replacements.items():
                    if old in p.text:
                        p.text = p.text.replace(old, new)

# Helper to find shape by text fragment
def find_shape_by_text(slide, text_fragment):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if text_fragment in shape.text_frame.text:
                return shape
    return None

# ==================== SLIDE 1 ====================
print("Processing Slide 1...")
slide1 = prs.slides[0]
slide1_replacements = {
    "Topic:": "Topic: INZIRA: A Machine Learning-Based Personalized Recommendation System for Rwanda",
    "Student Name:": "Student Name: Elissa Twizeyimana (BSc Software Engineering)",
    "Supervisor Name:": "Supervisor Name: Junior Turatsinze",
    "Department Name:": "Department Name: Software Engineering",
    "University Name:": "University Name: African Leadership University",
    "Presentation Date": "Presentation Date: July 2026"
}
replace_text_in_paragraphs(slide1, slide1_replacements)

# ==================== SLIDE 2 ====================
print("Processing Slide 2...")
slide2 = prs.slides[1]

# Clear description of the problem
shape_prob = find_shape_by_text(slide2, "Clear description of the problem")
if shape_prob:
    shape_prob.text_frame.clear()
    p = shape_prob.text_frame.paragraphs[0]
    p.text = "Existing recommendation platforms (like Google Maps & TripAdvisor) prioritize review volume, causing severe popularity bias. Additionally, traditional collaborative filtering algorithms suffer from cold-start failures, completely obscuring local hidden gems."
    p.font.name = "Arial"
    p.font.size = Pt(11.5)
    p.space_after = Pt(4)

# Brief context of the problem domain
shape_bg = find_shape_by_text(slide2, "Brief context of the problem domain")
if shape_bg:
    shape_bg.text_frame.clear()
    p = shape_bg.text_frame.paragraphs[0]
    p.text = "Tourism is Rwanda's leading economic driver, generating $647M in 2024. While major parks dominate tourist traffic, Rwanda's secondary cultural, historic, and community-based heritage sites remain digitally invisible."
    p.font.name = "Arial"
    p.font.size = Pt(11.5)
    p.space_after = Pt(4)

# Handle lists (Data/statistics and Stakeholders)
shape_stats = find_shape_by_text(slide2, "Data/statistics supporting the problem")
if shape_stats:
    shape_stats.text_frame.clear()
    stats = [
        "Evidence of Need",
        "• 85.6% of survey respondents report difficulty discovering niche cultural spots.",
        "• 95.0% demand personalized, location-aware travel recommendations.",
        "• Visitor traffic is heavily concentrated in Volcanoes/Akagera, leaving local sites under-visited."
    ]
    for idx, stat in enumerate(stats):
        p = shape_stats.text_frame.paragraphs[0] if idx == 0 else shape_stats.text_frame.add_paragraph()
        p.text = stat
        p.font.name = "Arial"
        if idx > 0:
            p.font.size = Pt(11)
            p.space_after = Pt(2)
        else:
            p.font.bold = True
            p.font.size = Pt(13)
            p.space_after = Pt(4)

shape_stakeholders = find_shape_by_text(slide2, "Target users and affected parties")
if shape_stakeholders:
    shape_stakeholders.text_frame.clear()
    stakeholders = [
        "Stakeholders",
        "• International Tourists: Seeking authentic heritage.",
        "• Local Residents: Looking for weekend leisure/dining.",
        "• Art & Craft Cooperatives: Seeking visitor visibility.",
        "• Regional Economy: Benefiting from decentralized tourism."
    ]
    for idx, sh in enumerate(stakeholders):
        p = shape_stakeholders.text_frame.paragraphs[0] if idx == 0 else shape_stakeholders.text_frame.add_paragraph()
        p.text = sh
        p.font.name = "Arial"
        if idx > 0:
            p.font.size = Pt(11)
            p.space_after = Pt(2)
        else:
            p.font.bold = True
            p.font.size = Pt(13)
            p.space_after = Pt(4)

# ==================== SLIDE 3 ====================
print("Processing Slide 3...")
slide3 = prs.slides[2]
slide3_replacements = {
    "[Overall project goal statement]": "Build a personalized and location-aware travel recommendation system (Inzira) for Rwanda that resolves popularity bias by surfacing local cultural, historic, and natural hidden gems, with full offline functionality."
}
replace_text_in_paragraphs(slide3, slide3_replacements)

# Objectives list
shape_objs = find_shape_by_text(slide3, "Objective 1")
if shape_objs:
    for p in shape_objs.text_frame.paragraphs:
        if "Objective 1" in p.text:
            p.text = "1. Curate and standardize a spatial catalog of 1,264 verified Rwandan points of interest (POIs)."
        elif "Objective 2" in p.text:
            p.text = "2. Train and serialize a RandomForestRegressor for personalized utility scoring."
        elif "Objective 3" in p.text:
            p.text = "3. Implement visitor-type weight adjustments and a responsive Next.js frontend with offline fallbacks."

# Scope list
shape_scope = find_shape_by_text(slide3, "Included:")
if shape_scope:
    shape_scope.text_frame.clear()
    scope_points = [
        "🔍 Scope",
        "Included:",
        "  • Stage 1 & 2 recommendations with dynamic Visitor Persona adjustments.",
        "  • Interactive Leaflet.js mapping with GPS routing links.",
        "  • Bookmarks folder storage (Google OAuth & local credentials).",
        "  • User feedback loop with 1-5 star ratings and comments.",
        "Excluded:",
        "  • Flight ticket bookings and reservation payment transactions.",
        "  • Native mobile app stores publishing (omitted for 8-week MVP timeline)."
    ]
    for idx, pt in enumerate(scope_points):
        p = shape_scope.text_frame.paragraphs[0] if idx == 0 else shape_scope.text_frame.add_paragraph()
        p.text = pt
        if idx == 0:
            p.font.bold = True
            p.font.size = Pt(16)
        elif pt.strip().startswith("Included:") or pt.strip().startswith("Excluded:"):
            p.font.bold = True
            p.font.size = Pt(13)
            p.space_before = Pt(6)
        else:
            p.font.size = Pt(11)
            p.space_after = Pt(2)

# ==================== SLIDE 4 ====================
print("Processing Slide 4...")
slide4 = prs.slides[3]
# Find the table shape
table_shape = None
for shape in slide4.shapes:
    if shape.has_table:
        table_shape = shape
        break

if table_shape:
    table = table_shape.table
    # Row 1 (System A)
    row1 = ["Official Portals", "Static brochures and directories", "No personalization or active recommendations", "Stage 1 custom preference utility vector match"]
    # Row 2 (System B)
    row2 = ["Google & TripAdvisor", "Crowdsourced review counts & global maps search", "Popularity bias (ignores small local spots)", "Hidden-gem scoring boost & hard popularity filtering"]
    
    for col_idx in range(4):
        table.cell(1, col_idx).text = row1[col_idx]
        table.cell(2, col_idx).text = row2[col_idx]
        
        # Style cell text
        for cell in [table.cell(1, col_idx), table.cell(2, col_idx)]:
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Arial"

# ==================== SLIDE 5 ====================
print("Processing Slide 5...")
slide5 = prs.slides[4]
shape_func = find_shape_by_text(slide5, "User registration and authentication")
if shape_func:
    for p in shape_func.text_frame.paragraphs:
        if "User registration" in p.text:
            p.text = "• Dynamic user onboarding (Tourist vs Resident, budget, time, interests)"
        elif "Data input and processing" in p.text:
            p.text = "• Stage 1 Main recommendations & Stage 2 spatial nearby search"
        elif "Reporting and analytics" in p.text:
            p.text = "• Interactive Leaflet mapping with external GPS redirection"
        elif "System administration" in p.text:
            p.text = "• Bookmarks folder storage (Google OAuth & local credentials)"
        elif "etc" in p.text:
            p.text = "• Star rating & comment feedback logging"

shape_nfunc = find_shape_by_text(slide5, "Performance: Response time")
if shape_nfunc:
    for p in shape_nfunc.text_frame.paragraphs:
        if "Performance:" in p.text:
            p.text = "• Performance: Server response & map rendering under 2.0s"
        elif "Security:" in p.text:
            p.text = "• Security: PBKDF2 (SHA-512) password hashing and Google OAuth integration"
        elif "Usability:" in p.text:
            p.text = "• Usability: Responsive layouts with offline client-side JS recommender fallback"
        elif "Scalability:" in p.text:
            p.text = "• Scalability: Schema supporting 1,264+ POIs & CSV feedback logs"

# Requirements gathering text box
shape_gath = find_shape_by_text(slide5, "Requirements Gathering:")
if shape_gath:
    for p in shape_gath.text_frame.paragraphs:
        if "Requirements Gathering:" in p.text:
            p.text = "Requirements Gathering: 120+ local travel surveys, TripAdvisor catalog analysis, and RDB reports."

# ==================== SLIDE 6 ====================
print("Processing Slide 6...")
slide6 = prs.slides[5]

# Change Class Diagram title to Two-Stage ML Flow
shape_class_title = find_shape_by_text(slide6, "Class Diagram")
if shape_class_title:
    for p in shape_class_title.text_frame.paragraphs:
        if "Class Diagram" in p.text:
            p.text = "📊  Two-Stage ML Flow"

# Helper to replace picture by coordinates
def replace_picture(slide, shape_index, new_image_path, scale_factor=1.0):
    shape = slide.shapes[shape_index]
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    
    # Calculate new sizes if scaled
    if scale_factor != 1.0:
        new_width = width * scale_factor
        new_height = height * scale_factor
        # Centering adjustment
        left = left - (new_width - width) / 2
        top = top - (new_height - height) / 2
        width = new_width
        height = new_height

    # Delete shape
    sp = shape._element
    sp.getparent().remove(sp)
    
    # Add new image
    slide.shapes.add_picture(new_image_path, left, top, width=width, height=height)

# Slide 6 Picture indices: Shape 2 (Architecture), Shape 4 (ERD), Shape 6 (Class Diagram placeholder)
replace_picture(slide6, 2, SYS_ARCH_PATH, scale_factor=1.0)
replace_picture(slide6, 4, ERD_PATH, scale_factor=1.0)
replace_picture(slide6, 6, ML_FLOW_PATH, scale_factor=1.0)

# ==================== SLIDE 7 ====================
print("Processing Slide 7...")
slide7 = prs.slides[6]
slide7_replacements = {
    "React, HTML, CSS": "Next.js (React), HSL CSS style configurations",
    "Node.js, Express, etc": "FastAPI (Python), Uvicorn backend server",
    "MySQL, MongoDB": "SQLite (inzira.db), JSON Catalog (1,264 POIs)",
}
replace_text_in_paragraphs(slide7, slide7_replacements)

shape_impl_details = find_shape_by_text(slide7, "RESTful API endpoints")
if shape_impl_details:
    for p in shape_impl_details.text_frame.paragraphs:
        if "RESTful API" in p.text:
            p.text = "• RESTful API endpoints (/api/recommend and /api/nearby) serving ML predictions"
        elif "Authentication using" in p.text:
            p.text = "• Dynamic Visitor Persona: Weight shifting for local residents vs tourists"
        elif "Real-time updates" in p.text:
            p.text = "• Resiliency Fallback: Replicated recommenders in client-side JS for 100% offline uptime"
        elif "Responsive design" in p.text:
            p.text = "• Secure User Directory: PBKDF2 (SHA-512) password hashing with random salt values to secure credentials"

# ==================== SLIDE 8 ====================
print("Processing Slide 8...")
slide8 = prs.slides[7]

shape_t1 = find_shape_by_text(slide8, "Individual components")
if shape_t1:
    shape_t1.text_frame.text = "• FastAPI endpoint check & JSON format schema parsed cleanly"
    shape_t1.text_frame.paragraphs[0].font.size = Pt(11)

shape_t2 = find_shape_by_text(slide8, "Components interactions")
if shape_t2:
    shape_t2.text_frame.text = "• Frontend-backend integration testing checking JSON body payload parameters"
    shape_t2.text_frame.paragraphs[0].font.size = Pt(11)

shape_t3 = find_shape_by_text(slide8, "End-to-End validation")
if shape_t3:
    shape_t3.text_frame.text = "• Complete UI walkthrough, offline fallback triggers, and catalog coordinate mapping checks"
    shape_t3.text_frame.paragraphs[0].font.size = Pt(11)

# Find and update slide 8 table
table_shape8 = None
for shape in slide8.shapes:
    if shape.has_table:
        table_shape8 = shape
        break

if table_shape8:
    table = table_shape8.table
    row1 = ["Model Fit (R²)", "≥ 0.95", "0.986", "✓ Passed"]
    row2 = ["Latency", "< 2.0s", "0.86s", "✓ Passed"]
    row3 = ["Coverage", "≥ 80%", "91.3%", "✓ Passed"]
    
    for col_idx in range(4):
        table.cell(1, col_idx).text = row1[col_idx]
        table.cell(2, col_idx).text = row2[col_idx]
        table.cell(3, col_idx).text = row3[col_idx]
        
        # Style cells
        for row_idx in [1, 2, 3]:
            for p in table.cell(row_idx, col_idx).text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Arial"

shape_val_method = find_shape_by_text(slide8, "Validation Method:")
if shape_val_method:
    shape_val_method.text_frame.text = "Validation Method: 5-Fold Cross-Validation, local user acceptance testing, E2E test runs"
    shape_val_method.text_frame.paragraphs[0].font.size = Pt(11)

# ==================== SLIDE 9 ====================
print("Processing Slide 9...")
slide9 = prs.slides[8]

shape_mo = find_shape_by_text(slide9, "Successfully developed and deployed")
if shape_mo:
    shape_mo.text_frame.clear()
    outcomes = [
        "• Developed and deployed a hybrid 2-Stage personalized and location-aware recommender system for Rwanda.",
        "• Standardized the first open-source spatial database of 1,264 points of interest."
    ]
    for idx, out in enumerate(outcomes):
        p = shape_mo.text_frame.paragraphs[0] if idx == 0 else shape_mo.text_frame.add_paragraph()
        p.text = out
        p.font.size = Pt(11)
        p.space_after = Pt(4)

shape_str = find_shape_by_text(slide9, "User-friendly interface")
if shape_str:
    shape_str.text_frame.clear()
    strengths = [
        "• High Prediction Quality: Precision@5 score of 0.820 (a +86.4% gain over baseline).",
        "• High Catalog Coverage: 91.3% coverage ensures small local spots get recommended.",
        "• Balanced Tourism: Diverting crowds from famous sites to local hidden gems.",
        "• Zero-Downtime Offline Mode: Seamless client-side JavaScript engine calculations."
    ]
    for idx, st in enumerate(strengths):
        p = shape_str.text_frame.paragraphs[0] if idx == 0 else shape_str.text_frame.add_paragraph()
        p.text = st
        p.font.size = Pt(11)
        p.space_after = Pt(2)

shape_lim = find_shape_by_text(slide9, "Limited to specific use cases")
if shape_lim:
    shape_lim.text_frame.clear()
    limitations = [
        "• Heavy dependency on accurate initial catalog tagging.",
        "• Cold-start for completely un-tagged new physical places.",
        "• Feedback loop logs currently rely on flat CSV files (requires DB scaling).",
        "• Web mapping requires active GPS coordinates for optimal nearby accuracy."
    ]
    for idx, lim in enumerate(limitations):
        p = shape_lim.text_frame.paragraphs[0] if idx == 0 else shape_lim.text_frame.add_paragraph()
        p.text = lim
        p.font.size = Pt(11)
        p.space_after = Pt(2)

shape_how = find_shape_by_text(slide9, "E.g The system successfully")
if shape_how:
    shape_how.text_frame.text = "Inzira effectively resolves the popularity-bias gap of Google Maps/TripAdvisor by boosting hidden gem visibility and dynamic traveler weighting. The system bridges the cold-start gap for niche community spots, driving local visitor engagement."
    shape_how.text_frame.paragraphs[0].font.size = Pt(11)

# ==================== SLIDE 10 ====================
print("Processing Slide 10...")
slide10 = prs.slides[9]

# Re-title Strength -> Key Conclusions, Limitation -> Future Work, How Results Address -> Contribution Statement
replace_text_in_paragraphs(slide10, {
    "✅ Strength": "🚀 Key Conclusions",
    "⚠️ Limitation": "🔮 Future Work",
    "💡How Results Address the Problem": "🌟 Contribution Statement"
})

shape10_mo = find_shape_by_text(slide10, "Successfully developed and deployed")
if shape10_mo:
    shape10_mo.text_frame.text = "Successfully proved that a hybrid content-based Random Forest model with visitor-type weighting can bypass the recommendation cold-start barrier and counter popularity bias."
    shape10_mo.text_frame.paragraphs[0].font.size = Pt(11)

shape10_str = find_shape_by_text(slide10, "User-friendly interface")
if shape10_str:
    shape10_str.text_frame.clear()
    conclusions = [
        "• Content-based constraints provide robust recommendations without dense historical rating matrices.",
        "• Hidden-gem filtering is a viable tool for spreading tourism benefits beyond major commercial centers.",
        "• Client-side fallback architectures ensure high application reliability in rural connectivity-sparse regions."
    ]
    for idx, c in enumerate(conclusions):
        p = shape10_str.text_frame.paragraphs[0] if idx == 0 else shape10_str.text_frame.add_paragraph()
        p.text = c
        p.font.size = Pt(11)
        p.space_after = Pt(2)

shape10_lim = find_shape_by_text(slide10, "Limited to specific use cases")
if shape10_lim:
    shape10_lim.text_frame.clear()
    future_works = [
        "• RDB Database Integration: Auto-sync with official Visit Rwanda directories.",
        "• Native Mobile Apps: Add geofenced push notifications and offline map caching.",
        "• Multi-Lingual Support: Localize to Kinyarwanda, French, and Swahili.",
        "• Dynamic Route Generation: Build multi-stop trip itineraries."
    ]
    for idx, f in enumerate(future_works):
        p = shape10_lim.text_frame.paragraphs[0] if idx == 0 else shape10_lim.text_frame.add_paragraph()
        p.text = f
        p.font.size = Pt(11)
        p.space_after = Pt(2)

shape10_how = find_shape_by_text(slide10, "E.g The system successfully")
if shape10_how:
    shape10_how.text_frame.text = "Inzira contributes a novel, offline-resilient, and popularity-bias-aware recommender framework tailored to emerging tourism economies, promoting local business discovery and preserving cultural heritage accessibility."
    shape10_how.text_frame.paragraphs[0].font.size = Pt(11)

# ==================== SLIDE 11 ====================
print("Processing Slide 11...")
slide11 = prs.slides[10]

# Set Title details
for shape in slide11.shapes:
    if shape.has_text_frame:
        if not shape.text_frame.text.strip():
            shape.text_frame.text = "INZIRA: A Machine Learning-Based Personalized Recommendation System for Rwanda"
            shape.text_frame.paragraphs[0].font.size = Pt(12)
            shape.text_frame.paragraphs[0].font.italic = True
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            break

# Save presentation
prs.save(OUTPUT_PATH)
print(f"\nPresentation populated from template and saved successfully to {OUTPUT_PATH}")
